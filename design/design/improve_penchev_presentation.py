from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


INPUT = Path(r"c:\Users\kospen\Documents\New project 6\Penchev_EDIH_AgriFood_10min_presentation.pptx")
OUTPUT = Path(r"d:\koceboce\KoceBoce\knowledge\design\Penchev_EDIH_AgriFood_10min_presentation_improved.pptx")
ASSET_DIR = Path(r"d:\koceboce\KoceBoce\knowledge\design\generated_assets")


COLORS = {
    "forest": RGBColor(67, 120, 67),
    "teal": RGBColor(46, 130, 126),
    "blue": RGBColor(65, 104, 144),
    "orange": RGBColor(224, 142, 65),
    "purple": RGBColor(110, 76, 147),
    "ink": RGBColor(32, 45, 43),
    "muted": RGBColor(88, 99, 95),
    "pale": RGBColor(242, 246, 240),
    "pale2": RGBColor(250, 252, 249),
    "line": RGBColor(210, 220, 214),
    "white": RGBColor(255, 255, 255),
}


def rgb_hex(rgb):
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def inches(value):
    return Inches(value)


def clear_slide(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=None, width=0.75):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(inches(0.72), inches(0.38), inches(8.7), inches(0.68))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = COLORS["ink"]
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS["muted"]
        p2.space_before = Pt(3)


def add_text(slide, x, y, w, h, lines, size=15, color=None, bold_first=False, line_spacing=1.05):
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inches(0.03)
    tf.margin_right = inches(0.03)
    tf.margin_top = inches(0.02)
    tf.margin_bottom = inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.font.bold = bold_first and i == 0
        p.space_after = Pt(5 if size >= 14 else 3)
        p.line_spacing = line_spacing
    return box


def add_label_value(slide, x, y, w, h, label, value, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    set_fill(card, COLORS["pale2"])
    set_line(card, COLORS["line"], 0.6)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x), inches(y), inches(0.08), inches(h))
    set_fill(strip, accent)
    set_line(strip)
    add_text(slide, x + 0.22, y + 0.12, w - 0.36, h - 0.18, [label, value], size=12, bold_first=True)


def add_stat_card(slide, x, y, w, h, number, label, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(w), inches(h))
    set_fill(card, COLORS["pale2"])
    set_line(card, COLORS["line"], 0.5)
    tf = card.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.font.name = "Calibri"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = "Calibri"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLORS["muted"]
    p2.alignment = PP_ALIGN.CENTER


def add_progress_bar(slide, x, y, label, value, max_value, color):
    add_text(slide, x, y - 0.03, 2.25, 0.28, [label], size=10.5, color=COLORS["ink"])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x + 2.45), inches(y + 0.03), inches(4.8), inches(0.18))
    set_fill(bg, RGBColor(230, 235, 232))
    set_line(bg)
    width = 4.8 * value / max_value
    fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inches(x + 2.45), inches(y + 0.03), inches(width), inches(0.18))
    set_fill(fg, color)
    set_line(fg)
    add_text(slide, x + 7.38, y - 0.04, 0.65, 0.3, [f"{value:g}%"], size=11, color=color)


def make_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    # Title illustration: farm fields connected to a digital hub.
    img = Image.new("RGBA", (900, 560), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((40, 340, 860, 520), radius=28, fill=(242, 246, 240, 255), outline=(210, 220, 214, 255), width=3)
    for i, color in enumerate([(67, 120, 67), (46, 130, 126), (224, 142, 65), (65, 104, 144)]):
        d.polygon([(70 + i * 195, 490), (235 + i * 195, 490), (190 + i * 195, 360), (100 + i * 195, 360)], fill=color + (210,))
    d.rounded_rectangle((350, 70, 550, 230), radius=30, fill=(250, 252, 249, 255), outline=(67, 120, 67, 255), width=6)
    d.ellipse((405, 104, 495, 194), fill=(67, 120, 67, 255))
    for x, y in [(150, 300), (300, 300), (600, 300), (750, 300), (450, 300)]:
        d.line((450, 230, x, y), fill=(65, 104, 144, 200), width=5)
        d.ellipse((x - 18, y - 18, x + 18, y + 18), fill=(250, 252, 249, 255), outline=(65, 104, 144, 255), width=5)
    path = ASSET_DIR / "agri_digital_network.png"
    img.save(path)

    # Country pattern abstract map.
    img = Image.new("RGBA", (900, 420), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    coords = {
        "Spain": (160, 260, (65, 104, 144)),
        "Slovenia": (410, 190, (46, 130, 126)),
        "Bulgaria": (545, 245, (67, 120, 67)),
        "Romania": (570, 170, (110, 76, 147)),
        "Greece": (610, 320, (67, 120, 67)),
        "Turkey": (725, 305, (224, 142, 65)),
    }
    d.rounded_rectangle((20, 20, 880, 400), radius=30, fill=(250, 252, 249, 255), outline=(210, 220, 214, 255), width=3)
    for x1, y1, x2, y2 in [(110, 250, 410, 190), (410, 190, 570, 170), (570, 170, 545, 245), (545, 245, 610, 320), (610, 320, 725, 305)]:
        d.line((x1, y1, x2, y2), fill=(165, 175, 165, 255), width=4)
    for name, (x, y, color) in coords.items():
        d.ellipse((x - 20, y - 20, x + 20, y + 20), fill=color + (255,), outline=(255, 255, 255, 255), width=4)
        d.text((x + 26, y - 12), name, fill=(32, 45, 43, 255))
    path = ASSET_DIR / "country_network.png"
    img.save(path)


def add_framework_matrix(slide):
    x0, y0, cell_w, cell_h = 0.82, 1.75, 2.02, 1.18
    items = [
        ("Economic viability", "ROI, cost and payback", COLORS["forest"]),
        ("Operational efficiency", "resource and workflow gains", COLORS["blue"]),
        ("Sustainability outcomes", "inputs, monitoring, resilience", COLORS["teal"]),
        ("Accessibility", "reach for smaller actors", COLORS["orange"]),
    ]
    for idx, (head, body, accent) in enumerate(items):
        x = x0 + (idx % 2) * (cell_w + 0.28)
        y = y0 + (idx // 2) * (cell_h + 0.28)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(cell_w), inches(cell_h))
        set_fill(card, COLORS["pale2"])
        set_line(card, COLORS["line"], 0.6)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x + 0.16), inches(y + 0.18), inches(0.34), inches(0.34))
        set_fill(circ, accent)
        set_line(circ)
        add_text(slide, x + 0.62, y + 0.12, cell_w - 0.75, 0.45, [head], size=12, color=COLORS["ink"], bold_first=True)
        add_text(slide, x + 0.18, y + 0.68, cell_w - 0.32, 0.32, [body], size=9.5, color=COLORS["muted"])


def build():
    make_assets()
    prs = Presentation(str(INPUT))
    for slide in prs.slides:
        clear_slide(slide)

    # Slide 1
    s = prs.slides[0]
    add_text(s, 0.78, 0.72, 5.65, 0.42, ["10-minute conference presentation"], size=13, color=COLORS["forest"], bold_first=True)
    add_text(
        s,
        0.78,
        1.18,
        5.95,
        1.78,
        ["Critical Analysis of European Digital Innovation Hubs in the Agri-Food Sector"],
        size=28,
        color=COLORS["ink"],
        bold_first=True,
        line_spacing=0.95,
    )
    add_text(
        s,
        0.82,
        3.95,
        5.25,
        0.82,
        ["Kostadin Penchev; Konstantin Stoyanov", "Trakia University, Stara Zagora, Bulgaria"],
        size=14,
        color=COLORS["muted"],
    )
    s.shapes.add_picture(str(ASSET_DIR / "agri_digital_network.png"), inches(5.75), inches(1.35), inches(3.6), inches(2.25))
    add_stat_card(s, 6.05, 4.36, 1.0, 0.72, "6", "country cases", COLORS["blue"])
    add_stat_card(s, 7.22, 4.36, 1.0, 0.72, "4", "evaluation criteria", COLORS["teal"])
    add_stat_card(s, 8.39, 4.36, 1.0, 0.72, "10", "minutes", COLORS["orange"])

    # Slide 2
    s = prs.slides[1]
    add_title(s, "Research Problem", "Why EDIH effectiveness matters for agri-food digitalization")
    add_text(
        s,
        0.78,
        1.32,
        5.05,
        2.35,
        [
            "Digital transformation is now a strategic priority in European agriculture.",
            "Yet adoption remains uneven, especially among small and medium-sized farms.",
            "The main barriers are not only technical: cost, uncertain returns, skills gaps and fragmented support limit uptake.",
        ],
        size=16,
    )
    add_label_value(s, 6.2, 1.38, 2.85, 0.9, "Adoption is socio-technical", "technology + economics + routines", COLORS["forest"])
    add_label_value(s, 6.2, 2.55, 2.85, 0.9, "Small farms need mediation", "testing, advice and trusted support", COLORS["teal"])
    add_label_value(s, 6.2, 3.72, 2.85, 0.9, "Assessment lens", "accessibility and inclusion", COLORS["orange"])
    add_text(s, 0.82, 4.62, 4.95, 0.65, ["Core question: do EDIHs create practical conditions for inclusive agri-food digital uptake?"], size=15, color=COLORS["forest"], bold_first=True)

    # Slide 3
    s = prs.slides[2]
    add_title(s, "Aim and Contribution")
    add_text(
        s,
        0.78,
        1.28,
        4.9,
        1.55,
        [
            "Aim: critically assess the potential of EDIHs to support the agri-food digital transition.",
            "The presentation compares Bulgaria, Romania, Slovenia, Greece, Spain and Turkey.",
        ],
        size=16,
    )
    add_label_value(s, 0.86, 3.15, 2.55, 0.92, "Contribution 1", "moves from description to evaluation", COLORS["blue"])
    add_label_value(s, 3.65, 3.15, 2.55, 0.92, "Contribution 2", "links EDIHs to rural inclusion", COLORS["teal"])
    add_label_value(s, 6.44, 3.15, 2.55, 0.92, "Contribution 3", "compares different ecosystem conditions", COLORS["orange"])
    for i, country in enumerate(["Bulgaria", "Romania", "Slovenia", "Greece", "Spain", "Turkey"]):
        x = 0.9 + (i % 3) * 1.72
        y = 4.55 + (i // 3) * 0.46
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inches(x), inches(y), inches(1.35), inches(0.34))
        set_fill(chip, COLORS["pale"])
        set_line(chip, COLORS["line"], 0.4)
        chip.text_frame.text = country
        p = chip.text_frame.paragraphs[0]
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLORS["ink"]
        p.alignment = PP_ALIGN.CENTER
    s.shapes.add_picture(str(ASSET_DIR / "country_network.png"), inches(5.85), inches(1.15), inches(3.35), inches(1.56))

    # Slide 4
    s = prs.slides[3]
    add_title(s, "Materials and Methods", "Qualitative comparative design")
    add_text(
        s,
        0.78,
        1.25,
        4.3,
        1.25,
        [
            "Evidence base combines peer-reviewed literature with institutional and policy documents.",
            "The six cases were selected to capture different agri-food innovation ecosystems.",
        ],
        size=15.5,
    )
    steps = [
        ("1", "Review", "scientific literature on EDIHs, agriculture and digital adoption", COLORS["blue"]),
        ("2", "Document analysis", "EU, national and hub-level sources", COLORS["forest"]),
        ("3", "Comparative interpretation", "country patterns and structural constraints", COLORS["teal"]),
    ]
    for i, (num, head, body, accent) in enumerate(steps):
        x = 0.9 + i * 2.95
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, inches(x), inches(3.02), inches(0.48), inches(0.48))
        set_fill(circ, accent)
        set_line(circ)
        circ.text_frame.text = num
        circ.text_frame.paragraphs[0].font.bold = True
        circ.text_frame.paragraphs[0].font.size = Pt(13)
        circ.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]
        circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text(s, x + 0.6, 2.94, 2.1, 0.86, [head, body], size=11.5, bold_first=True)
        if i < 2:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(x + 2.42), inches(3.14), inches(0.35), inches(0.2))
            set_fill(arr, RGBColor(165, 175, 165))
            set_line(arr)
    add_label_value(s, 5.7, 1.25, 3.35, 1.05, "Scope", "Bulgaria, Romania, Slovenia, Greece, Spain and Turkey", COLORS["purple"])

    # Slide 5
    s = prs.slides[4]
    add_title(s, "Analytical Framework")
    add_text(
        s,
        5.4,
        1.42,
        3.65,
        1.35,
        [
            "Four criteria structure the comparison.",
            "Barriers are treated as contextual conditions that shape whether services become usable.",
        ],
        size=15,
    )
    add_framework_matrix(s)
    add_label_value(s, 5.35, 3.35, 3.65, 1.35, "Contextual barriers", "digital literacy, broadband connectivity, finance, fragmented holdings", COLORS["purple"])

    # Slide 6
    s = prs.slides[5]
    add_title(s, "Network Reach vs. Agricultural Uptake", "Percentages and counts reported in the article")
    add_text(
        s,
        0.78,
        1.18,
        8.55,
        0.72,
        ["The EDIH network is broad, but agriculture appears under-represented among reported beneficiaries."],
        size=15.5,
        color=COLORS["ink"],
    )
    add_progress_bar(s, 0.95, 2.15, "Manufacturing beneficiaries", 28, 40, COLORS["blue"])
    add_progress_bar(s, 0.95, 2.78, "Agricultural users", 1.7, 40, COLORS["forest"])
    add_progress_bar(s, 0.95, 3.41, "Agri-food-related hubs", 36, 40, COLORS["teal"])
    add_progress_bar(s, 0.95, 4.04, "Exclusive agricultural mandate", 5, 40, COLORS["orange"])
    add_stat_card(s, 0.98, 5.04, 1.55, 0.7, "254", "EDIHs by May 2025", COLORS["blue"])
    add_stat_card(s, 3.02, 5.04, 1.55, 0.7, "93", "agri-food support", COLORS["teal"])
    add_stat_card(s, 5.06, 5.04, 1.55, 0.7, "12", "exclusive mandate", COLORS["orange"])
    add_stat_card(s, 7.1, 5.04, 1.55, 0.7, ">18k", "service interactions", COLORS["purple"])

    # Slide 7
    s = prs.slides[6]
    add_title(s, "How EDIHs Can Add Value")
    add_text(
        s,
        0.78,
        1.18,
        8.45,
        0.65,
        ["The strongest value proposition is not technology promotion alone, but reducing risk and translating digital tools into farm-level practice."],
        size=15.2,
    )
    chain = [
        ("Uncertainty", "context-specific returns", COLORS["blue"]),
        ("Test before invest", "realistic farm conditions", COLORS["forest"]),
        ("Training", "skills and routines", COLORS["teal"]),
        ("Finance & networks", "implementation pathways", COLORS["orange"]),
    ]
    for i, (head, body, accent) in enumerate(chain):
        x = 0.72 + i * 2.27
        add_label_value(s, x, 2.6, 1.95, 1.05, head, body, accent)
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(x + 1.98), inches(2.98), inches(0.34), inches(0.22))
            set_fill(arr, RGBColor(165, 175, 165))
            set_line(arr)
    add_label_value(s, 1.15, 4.58, 7.7, 0.82, "Critical qualification", "EDIHs create conditions under which benefits become more plausible; they do not guarantee outcomes.", COLORS["purple"])

    # Slide 8
    s = prs.slides[7]
    add_title(s, "Comparative Country Patterns")
    countries = [
        ("Bulgaria", "explicit agri-food positioning", COLORS["forest"]),
        ("Greece", "sector-specialized hub", COLORS["forest"]),
        ("Slovenia", "dedicated and multi-sector structures", COLORS["teal"]),
        ("Spain", "mature ecosystem and coordination platform", COLORS["blue"]),
        ("Turkey", "emerging agri-food innovation ecosystem", COLORS["orange"]),
        ("Romania", "heterogeneous multi-sector configuration", COLORS["purple"]),
    ]
    for i, (country, pattern, accent) in enumerate(countries):
        x = 0.72 + (i % 2) * 4.45
        y = 1.3 + (i // 2) * 1.05
        add_label_value(s, x, y, 3.95, 0.82, country, pattern, accent)
    s.shapes.add_picture(str(ASSET_DIR / "country_network.png"), inches(1.75), inches(4.72), inches(6.4), inches(1.5))

    # Slide 9
    s = prs.slides[8]
    add_title(s, "Discussion and Conclusions")
    add_text(
        s,
        0.78,
        1.18,
        8.55,
        0.85,
        ["EDIHs are potentially valuable intermediary structures, but their effectiveness is conditional rather than automatic."],
        size=16,
        color=COLORS["ink"],
        bold_first=True,
    )
    add_label_value(s, 0.84, 2.38, 2.6, 1.35, "1. Sector adaptation", "services must become concrete agricultural use cases", COLORS["forest"])
    add_label_value(s, 3.7, 2.38, 2.6, 1.35, "2. Inclusive access", "formal availability does not ensure uptake by small farms", COLORS["orange"])
    add_label_value(s, 6.56, 2.38, 2.6, 1.35, "3. Governance alignment", "EU, national and regional levels need closer coordination", COLORS["blue"])
    add_text(
        s,
        1.0,
        4.55,
        8.0,
        0.78,
        ["Future research should examine user experience, post-adoption outcomes and measurable farm-level effects."],
        size=15,
        color=COLORS["purple"],
        bold_first=True,
    )

    # Slide 10
    s = prs.slides[9]
    add_title(s, "Final Message")
    add_text(
        s,
        0.92,
        1.35,
        8.0,
        1.4,
        [
            "EDIHs can support agri-food digitalization when services are sector-specific, affordable and locally embedded.",
            "Inclusive impact depends on practical accessibility for smaller agricultural actors.",
        ],
        size=19,
        color=COLORS["ink"],
    )
    add_label_value(s, 1.1, 3.55, 7.8, 1.12, "Take-home message", "Relevance depends on specialization, ecosystem capacity and real accessibility.", COLORS["forest"])
    add_text(s, 1.12, 5.08, 7.7, 0.42, ["Thank you for your attention"], size=20, color=COLORS["forest"], bold_first=True)

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build()
